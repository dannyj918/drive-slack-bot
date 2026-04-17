"""
RAG indexer
===========
Standalone script that indexes Google Shared Drive files into a local Chroma
vector database for semantic search.

First run  → full sync: all files are indexed, a Changes API page token is saved.
Later runs → incremental sync: only files added/modified/deleted since the last
             run are re-indexed. This keeps re-runs fast.

Usage:
  python rag_indexer.py          # auto-detects full vs incremental
  python rag_indexer.py --full   # force a full re-index

Schedule with cron (every 30 minutes):
  */30 * * * * cd /path/to/drive-slack-bot && .venv/bin/python rag_indexer.py
"""

import argparse
import io
import logging
import os
import sys
import tempfile
import time

import anthropic
import base64
import chromadb
import docx
import openai
import openpyxl
import pypdf
from bs4 import BeautifulSoup
from pptx import Presentation
from dotenv import load_dotenv
from google.oauth2 import service_account
from googleapiclient.discovery import build
from googleapiclient.errors import HttpError
from googleapiclient.http import MediaIoBaseDownload

load_dotenv()

logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s [%(levelname)s] %(name)s: %(message)s",
)
logging.getLogger("pypdf").setLevel(logging.ERROR)
logger = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# Constants
# ---------------------------------------------------------------------------

EMBEDDING_MODEL    = "text-embedding-3-small"
COLLECTION_NAME    = "drive_docs"
CHANGES_TOKEN_FILE = "changes_token.txt"
CHUNK_SIZE         = 500   # words per chunk
CHUNK_OVERLAP      = 50    # word overlap between chunks
EMBED_BATCH_SIZE   = 100   # max texts per OpenAI embedding request

_SCOPES = ["https://www.googleapis.com/auth/drive.readonly"]

# MIME types that Drive can export as plain text
_EXPORTABLE: dict[str, str] = {
    "application/vnd.google-apps.document":     "text/plain",
    "application/vnd.google-apps.spreadsheet":  "text/csv",
    "application/vnd.google-apps.presentation": "text/plain",
    "application/vnd.google-apps.form":         "text/plain",
}

# ---------------------------------------------------------------------------
# Clients (lazy singletons)
# ---------------------------------------------------------------------------

_drive_service = None
_openai_client: openai.OpenAI | None = None


def _build_drive_service():
    global _drive_service
    if _drive_service is not None:
        return _drive_service
    key_file = os.environ.get("GOOGLE_SERVICE_ACCOUNT_FILE", "service_account.json")
    creds = service_account.Credentials.from_service_account_file(
        key_file, scopes=_SCOPES
    )
    _drive_service = build("drive", "v3", credentials=creds, cache_discovery=False)
    return _drive_service


def _get_openai_client() -> openai.OpenAI:
    global _openai_client
    if _openai_client is None:
        _openai_client = openai.OpenAI(api_key=os.environ["OPENAI_API_KEY"])
    return _openai_client


def _get_chroma_collection() -> chromadb.Collection:
    db_path = os.environ.get("CHROMA_DB_PATH", "./chroma_db")
    client = chromadb.PersistentClient(path=db_path)
    return client.get_or_create_collection(
        name=COLLECTION_NAME,
        metadata={"hnsw:space": "cosine"},
    )


# ---------------------------------------------------------------------------
# Drive helpers
# ---------------------------------------------------------------------------

def list_all_files() -> list[dict]:
    """Page through all non-trashed files in the Shared Drive."""
    drive_id = os.environ.get("SHARED_DRIVE_ID", "").strip()
    if not drive_id:
        raise ValueError("SHARED_DRIVE_ID is not set in .env")

    service = _build_drive_service()
    files = []
    page_token = None

    while True:
        params = dict(
            q="trashed = false",
            corpora="drive",
            driveId=drive_id,
            includeItemsFromAllDrives=True,
            supportsAllDrives=True,
            pageSize=100,
            fields="nextPageToken,files(id,name,mimeType,webViewLink,modifiedTime)",
        )
        if page_token:
            params["pageToken"] = page_token

        response = service.files().list(**params).execute()
        files.extend(response.get("files", []))
        page_token = response.get("nextPageToken")
        if not page_token:
            break

    logger.info("Found %d files in Shared Drive", len(files))
    return files


def extract_text(file: dict) -> str:
    """
    Extract plain text from a Drive file.

    - Google Workspace files (Docs, Sheets, Slides, Forms): exported via Drive API.
    - PDFs: downloaded and parsed with pypdf.
    - Everything else: skipped (returns "").
    """
    service   = _build_drive_service()
    file_id   = file["id"]
    mime_type = file.get("mimeType", "")
    name      = file.get("name", file_id)

    if mime_type in _EXPORTABLE:
        try:
            data = service.files().export(
                fileId=file_id,
                mimeType=_EXPORTABLE[mime_type],
            ).execute()
            return data.decode("utf-8", errors="ignore") if isinstance(data, bytes) else str(data)
        except HttpError as exc:
            logger.warning("Export failed for %r: %s", name, exc)
            return ""

    if mime_type == "application/pdf":
        try:
            request  = service.files().get_media(fileId=file_id)
            buf      = io.BytesIO()
            downloader = MediaIoBaseDownload(buf, request)
            done = False
            while not done:
                _, done = downloader.next_chunk()
            buf.seek(0)
            reader = pypdf.PdfReader(buf)
            text = "\n".join(
                (page.extract_text() or "") for page in reader.pages
            )
            return text
        except Exception as exc:
            logger.warning("PDF parse failed for %r: %s — trying Claude fallback", name, exc)
            try:
                buf.seek(0)
                pdf_b64 = base64.standard_b64encode(buf.read()).decode()
                ai = anthropic.Anthropic(api_key=os.environ["ANTHROPIC_API_KEY"])
                msg = ai.messages.create(
                    model=os.environ.get("CLAUDE_MODEL", "claude-sonnet-4-6"),
                    max_tokens=8096,
                    messages=[{
                        "role": "user",
                        "content": [
                            {
                                "type": "document",
                                "source": {"type": "base64", "media_type": "application/pdf", "data": pdf_b64},
                            },
                            {"type": "text", "text": "Extract and return all text content from this PDF. Return only the extracted text with no commentary."},
                        ],
                    }],
                )
                text = msg.content[0].text if msg.content else ""
                if text:
                    logger.info("Claude fallback succeeded for %r", name)
                return text
            except Exception as fb_exc:
                logger.warning("Claude fallback also failed for %r: %s", name, fb_exc)
                return ""

    # Microsoft Office formats uploaded to Drive
    _OFFICE_TYPES = {
        "application/vnd.openxmlformats-officedocument.presentationml.presentation": "pptx",
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document":   "docx",
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":         "xlsx",
    }
    if mime_type in _OFFICE_TYPES:
        try:
            request = service.files().get_media(fileId=file_id)
            buf = io.BytesIO()
            downloader = MediaIoBaseDownload(buf, request)
            done = False
            while not done:
                _, done = downloader.next_chunk()
            buf.seek(0)
            fmt = _OFFICE_TYPES[mime_type]
            if fmt == "pptx":
                prs = Presentation(buf)
                parts = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            parts.append(shape.text_frame.text)
                return "\n".join(parts)
            if fmt == "docx":
                doc = docx.Document(buf)
                return "\n".join(p.text for p in doc.paragraphs)
            if fmt == "xlsx":
                wb = openpyxl.load_workbook(buf, read_only=True, data_only=True)
                rows = []
                for sheet in wb.worksheets:
                    for row in sheet.iter_rows(values_only=True):
                        line = "\t".join("" if v is None else str(v) for v in row)
                        if line.strip():
                            rows.append(line)
                return "\n".join(rows)
        except Exception as exc:
            logger.warning("Office parse failed for %r: %s", name, exc)
            return ""

    if mime_type in ("text/html", "application/xhtml+xml"):
        try:
            request = service.files().get_media(fileId=file_id)
            buf = io.BytesIO()
            downloader = MediaIoBaseDownload(buf, request)
            done = False
            while not done:
                _, done = downloader.next_chunk()
            buf.seek(0)
            soup = BeautifulSoup(buf.read(), "html.parser")
            for tag in soup(["style", "noscript"]):
                tag.decompose()
            return soup.get_text(separator=" ", strip=True)
        except Exception as exc:
            logger.warning("HTML parse failed for %r: %s", name, exc)
            return ""

    logger.debug("Skipping unsupported file type %s for %r", mime_type, name)
    return ""


# ---------------------------------------------------------------------------
# Chunking + embedding
# ---------------------------------------------------------------------------

def chunk_text(text: str) -> list[str]:
    """Split text into overlapping word-count chunks."""
    words = text.split()
    if not words:
        return []
    chunks = []
    start = 0
    while start < len(words):
        end   = min(start + CHUNK_SIZE, len(words))
        chunk = " ".join(words[start:end]).strip()
        if chunk:
            chunks.append(chunk)
        start += CHUNK_SIZE - CHUNK_OVERLAP
    return chunks


def embed_texts(texts: list[str]) -> list[list[float]]:
    """Embed a list of strings using OpenAI text-embedding-3-small, in batches."""
    client = _get_openai_client()
    embeddings: list[list[float]] = []

    for i in range(0, len(texts), EMBED_BATCH_SIZE):
        batch = texts[i : i + EMBED_BATCH_SIZE]
        max_retries, delay = 3, 2.0
        for attempt in range(max_retries):
            try:
                response = client.embeddings.create(model=EMBEDDING_MODEL, input=batch)
                break
            except openai.RateLimitError:
                if attempt < max_retries - 1:
                    logger.warning(
                        "OpenAI rate limit hit — retrying in %.0fs (attempt %d/%d)",
                        delay, attempt + 1, max_retries,
                    )
                    time.sleep(delay)
                    delay *= 2
                else:
                    logger.error("OpenAI rate limit exceeded after %d retries", max_retries)
                    raise
        embeddings.extend(item.embedding for item in response.data)

    return embeddings


# ---------------------------------------------------------------------------
# Indexing
# ---------------------------------------------------------------------------

def index_file(file: dict, collection: chromadb.Collection) -> int:
    """
    Index a single Drive file into Chroma.
    Deletes any existing chunks for this file before inserting new ones.
    Returns the number of chunks stored (0 if the file was skipped).
    """
    file_id = file["id"]
    name    = file.get("name", file_id)

    # Remove stale chunks first — runs even if new text is empty, preventing orphans
    try:
        collection.delete(where={"file_id": file_id})
    except Exception:
        pass  # collection may be empty; delete(where=...) can raise if no matches

    text = extract_text(file)
    if not text.strip():
        logger.debug("No text extracted from %r — skipping", name)
        return 0

    chunks = chunk_text(text)
    if not chunks:
        return 0

    embeddings = embed_texts(chunks)
    ids        = [f"{file_id}_{i}" for i in range(len(chunks))]
    metadatas  = [
        {
            "file_id":      file_id,
            "file_name":    name,
            "file_link":    file.get("webViewLink", ""),
            "mime_type":    file.get("mimeType", ""),
            "chunk_index":  i,
            "modified_time": (file.get("modifiedTime") or "")[:10],
        }
        for i in range(len(chunks))
    ]

    collection.upsert(
        ids=ids,
        embeddings=embeddings,
        documents=chunks,
        metadatas=metadatas,
    )

    logger.info("Indexed %r → %d chunk(s)", name, len(chunks))
    return len(chunks)


# ---------------------------------------------------------------------------
# Changes API token persistence
# ---------------------------------------------------------------------------

def _save_token(token: str) -> None:
    with open(CHANGES_TOKEN_FILE, "w") as fh:
        fh.write(token)


def _load_token() -> str | None:
    if os.path.exists(CHANGES_TOKEN_FILE):
        with open(CHANGES_TOKEN_FILE) as fh:
            return fh.read().strip() or None
    return None


def _get_start_token() -> str:
    """Fetch the current Changes API page token (snapshot of Drive right now)."""
    drive_id = os.environ.get("SHARED_DRIVE_ID", "").strip()
    service  = _build_drive_service()
    response = service.changes().getStartPageToken(
        supportsAllDrives=True,
        driveId=drive_id,
    ).execute()
    return response["startPageToken"]


# ---------------------------------------------------------------------------
# Sync modes
# ---------------------------------------------------------------------------

def full_sync(collection: chromadb.Collection) -> None:
    """Index every file in the Shared Drive and save a Changes API token."""
    # Grab the token BEFORE listing files so we don't miss changes during indexing
    start_token = _get_start_token()

    files = list_all_files()
    total_chunks = 0

    for file in files:
        try:
            total_chunks += index_file(file, collection)
        except Exception as exc:
            logger.error("Failed to index %r: %s", file.get("name"), exc)

    _save_token(start_token)
    logger.info(
        "Full sync complete: %d file(s), %d chunk(s) indexed. Token saved.",
        len(files),
        total_chunks,
    )


def incremental_sync(collection: chromadb.Collection) -> None:
    """Re-index only files that changed since the last run (via Drive Changes API)."""
    token = _load_token()
    if not token:
        logger.warning("No changes token found — falling back to full sync")
        full_sync(collection)
        return

    drive_id = os.environ.get("SHARED_DRIVE_ID", "").strip()
    service  = _build_drive_service()
    processed = 0

    logger.info("Incremental sync from token: %s", token)

    # The Changes API may paginate; iterate until newStartPageToken is present
    while True:
        try:
            response = service.changes().list(
                pageToken=token,
                driveId=drive_id,
                includeItemsFromAllDrives=True,
                supportsAllDrives=True,
                fields=(
                    "nextPageToken,newStartPageToken,"
                    "changes(changeType,fileId,removed,"
                    "file(id,name,mimeType,webViewLink,modifiedTime,trashed))"
                ),
            ).execute()
        except HttpError as exc:
            if exc.resp.status == 400 or "invalid" in str(exc).lower():
                logger.warning(
                    "Changes API token is expired or invalid (HTTP %s) — "
                    "deleting token and falling back to full sync.",
                    exc.resp.status,
                )
                if os.path.exists(CHANGES_TOKEN_FILE):
                    os.remove(CHANGES_TOKEN_FILE)
                full_sync(collection)
                return
            raise

        for change in response.get("changes", []):
            if "fileId" not in change:
                continue
            file_id   = change["fileId"]
            file_meta = change.get("file") or {}
            removed   = change.get("removed") or file_meta.get("trashed", False)

            if removed:
                try:
                    collection.delete(where={"file_id": file_id})
                    logger.info("Removed deleted file %s from index", file_id)
                except Exception:
                    pass
            elif file_meta:
                try:
                    incoming_modified = (file_meta.get("modifiedTime") or "")[:10]
                    if incoming_modified:
                        existing = collection.get(
                            where={"file_id": file_id},
                            limit=1,
                            include=["metadatas"],
                        )
                        if (
                            existing["metadatas"]
                            and existing["metadatas"][0].get("modified_time") == incoming_modified
                        ):
                            logger.debug(
                                "Skipping %r — modifiedTime unchanged (%s)",
                                file_meta.get("name"), incoming_modified,
                            )
                            processed += 1
                            continue
                    index_file(file_meta, collection)
                except Exception as exc:
                    logger.error("Failed to re-index %r: %s", file_meta.get("name"), exc)
            else:
                logger.warning(
                    "Skipping change for file_id=%s: no file metadata returned",
                    file_id,
                )

            processed += 1

        new_token = response.get("newStartPageToken") or response.get("nextPageToken")
        if new_token:
            _save_token(new_token)

        # newStartPageToken signals the end of the change list
        if response.get("newStartPageToken"):
            break
        token = response.get("nextPageToken")
        if not token:
            break

    logger.info("Incremental sync complete: %d change(s) processed", processed)


# ---------------------------------------------------------------------------
# Entry point
# ---------------------------------------------------------------------------

def print_stats() -> None:
    db_path = os.environ.get("CHROMA_DB_PATH", "./chroma_db")
    if not os.path.exists(db_path):
        print(f"No Chroma DB found at {db_path} — run rag_indexer.py to build the index")
        return

    collection = _get_chroma_collection()
    chunk_count = collection.count()

    total_bytes = sum(
        os.path.getsize(os.path.join(root, f))
        for root, _, files in os.walk(db_path)
        for f in files
    )
    size_mb = total_bytes / (1024 * 1024)

    print(f"Chunks : {chunk_count:,}")
    print(f"DB size: {size_mb:.1f} MB  ({db_path})")


if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Index Shared Drive files into Chroma")
    parser.add_argument(
        "--full",
        action="store_true",
        help="Force a full re-index even if a changes token exists",
    )
    parser.add_argument(
        "--stats",
        action="store_true",
        help="Print vector DB stats (chunk count, disk size) and exit",
    )
    args = parser.parse_args()

    if args.stats:
        print_stats()
        sys.exit(0)

    collection = _get_chroma_collection()

    if args.full or not os.path.exists(CHANGES_TOKEN_FILE):
        logger.info("Starting full sync…")
        full_sync(collection)
    else:
        logger.info("Starting incremental sync…")
        incremental_sync(collection)
